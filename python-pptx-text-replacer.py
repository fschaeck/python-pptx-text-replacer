import sys
import argparse

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.shapes.graphfrm import GraphicFrame
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

"""
This module implements text replacement in Powerpoint files in pptx format.
The text is searched and replaced in all possible places.
"""
 
class python_pptx_text_replacer:
    """
    This class implements text replacement in Powerpoint files in pptx format.
    The text is searched and replaced in all possible places.
    """
 
    def __init__(self, presentation_file_name):
        self._replacements = []
        self._collected_replacements = []
        self._presentation_file_name = presentation_file_name
        self._presentation = Presentation(presentation_file_name)

    def replace_text(self, replacements):
        print(self._replacements)
        self._replacements = list( (self._ensure_unicode(match),self._ensure_unicode(repl)) for (match,repl) in replacements )
        self._collected_replacements.append(replacements)
        # loop through all slides
        slide_idx = 0
        print("Presentation[%s]" % (self._presentation_file_name))
        for slide in self._presentation.slides:
            print("  Slide[%s, id=%s] with title '%s'" % ( slide_idx, slide.slide_id, "<no title>" if slide.shapes.title is None else slide.shapes.title.text ))
            self._process_shapes(2, slide)
            slide_idx += 1

    def write_presentation_to_file(self, presentation_output_file_name):
        self._presentation.save(presentation_output_file_name)

    def get_replacements(self):
        return self._collected_replacements

    def get_presentation_file_name(self):
        return self._presentation_file_name

    def _ensure_unicode(self, text):
        if isinstance(text,(str,bytes) if sys.version_info.major==2 else bytes):
            return text.decode('UTF-8')
        return text

    def _replace_text_in_text_frame(self, level, text_frame):
        for (match, replacement) in self._replacements:
            pos_in_text_frame = self._ensure_unicode(text_frame.text).find(match)
            if pos_in_text_frame < 0:
                print("%sTrying to match '%s' -> no match" % ( "  "*level, match ))
            while pos_in_text_frame>=0:
                print("%sTrying to match '%s' -> matched at %s" % ( "  "*level, match, pos_in_text_frame ))
                to_match = match
                to_replace = replacement
                paragraph_idx = 0
                for paragraph in text_frame.paragraphs:
                    paragraph_len = len(self._ensure_unicode(paragraph.text))
                    if pos_in_text_frame >= paragraph_len:
                        pos_in_text_frame -= paragraph_len+1 # +1 for the new-line-character
                    else:
                        # this is the paragraph that contains the beginning of the match
                        (to_match, to_replace) = self._replace_runs_text(level+1, paragraph_idx, paragraph.runs, pos_in_text_frame, to_match, to_replace)
                        if len(to_match) == 0: # are we done with this match
                            break;
                        pos_in_text_frame = 0
                    paragraph_idx += 1
                pos_in_text_frame = self._ensure_unicode(text_frame.text).find(match)

    def _replace_runs_text(self, level, paragraph_idx, runs, pos, match, replacement):
        cnt = len(runs)
        i = 0
        while i<cnt:
            olen = len(self._ensure_unicode(runs[i].text))
            if pos>=olen:
                pos -= olen # the relative position of our match in the next run's text
                i += 1      # and off to the next run
            else:
                # we found the run, where the match starts!
                to_match = match
                match_len = len(to_match)
                to_replace = replacement
                repl_len = len(to_replace)

                while i<cnt:
                    run = runs[i]
                    otext = self._ensure_unicode(run.text)
                    olen = len(otext)
                    if pos+match_len < olen:
                        # our match ends before the end of the text of this run therefore
                        # we put the rest of our replacement string here and we are done!
                        run.text = otext[0:pos]+to_replace+otext[pos+match_len:]
                        print("%sRun[%s,%s]: '%s' -> '%s'" % ( "  "*level, paragraph_idx, i, otext, run.text ))
                        return ('','')
                    if pos+match_len == olen:
                        # our match ends together with the text of this run therefore
                        # we put the rest of our replacement string here and we are done!
                        run.text = otext[0:pos]+to_replace
                        print("%sRun[%s,%s]: '%s' -> '%s'" % ( "  "*level, paragraph_idx, i, otext, run.text ))
                        return ('','')
                    # we still haven't found all of our original match string
                    # so we process what we have here and go on to the next run
                    part_match_len = olen-pos
                    ntext = otext[0:pos]
                    if repl_len <= part_match_len:
                        # we now found at least as many characters for our match string
                        # as we have replacement characters for it. Thus we use up the
                        # the rest of our replacement string here and will replace the
                        # remainder of the match with an empty string (which happens
                        # to happen in this exact same spot for the next run ;-))
                        ntext += to_replace
                        repl_len = 0
                        to_replace = ''
                    else:
                        # we have got some more match characters but still more
                        # replacement characters than match characters found 
                        ntext += to_replace[0:part_match_len]
                        to_replace = to_replace[part_match_len:]
                        repl_len -= part_match_len
                    print("%sRun[%s,%s]: '%s' -> '%s'" % ( "  "*level, paragraph_idx, i, otext, ntext ))
                    run.text = ntext            # save the new text to the run
                    to_match = to_match[part_match_len:] # this is what is left to match
                    match_len -= part_match_len # this is the length of the match that is left
                    pos = 0                     # in the next run, we start at pos 0 with our match
                    i += 1                      # and off to the next run
                return (to_match, to_replace)
            

    def _process_text_frame(self, level, text_frame):
        print("%sTextFrame: '%s'" % ( "  "*level, text_frame.text ))
        paragraph_idx = 0
        for paragraph in text_frame.paragraphs:
            print("%sParagraph[%s]: '%s'" % ( "  "*(level+1), paragraph_idx, paragraph.text ))
            run_idx = 0
            for run in paragraph.runs:
                print("%sRun[%s,%s]: '%s'" % ( "  "*(level+2), paragraph_idx, run_idx, run.text ))
                run_idx += 1
            paragraph_idx += 1
        self._replace_text_in_text_frame(level+1,text_frame)

    def _process_shapes(self, level, shape_list_parent):
        for shape in shape_list_parent.shapes:
            print("%sShape[%s, id=%s, type=%s]" % ( "  "*level, shape_list_parent.shapes.index(shape), shape.shape_id, shape.shape_type ))
            if shape.has_text_frame:
                self._process_text_frame(level+1,shape.text_frame)
            if shape.has_table:
                table = shape.table
                row_cnt = len(table.rows)
                col_cnt = len(table.columns)
                print("%sTable[%s,%s]" % ( "  "*(level+1), row_cnt, col_cnt ) )
                for row in range(0, row_cnt):
                    for col in range(0, col_cnt):
                        cell = table.cell(row,col)
                        print("%sCell[%s,%s]: '%s'" % ( "  "*(level+2), row, col, cell.text ))
                        self._process_text_frame(level+3, cell.text_frame)
            if shape.shape_type==MSO_SHAPE_TYPE.GROUP:
                self._process_shapes(level+1, shape)

if __name__ == '__main__':
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument('--match',   '-m', action='append', required=True, dest='matches', metavar='<match>', help='the match to look for')
    p.add_argument('--replace', '-r', action='append', required=True, dest='replacements', metavar='<replacement>', help="the matches' replacement")
    p.add_argument('--input',   '-i', action='store',  required=True, metavar='<input file>', help="the file to replace the text in")
    p.add_argument('--output',  '-o', action='store',  required=True, metavar='<output file', help="the file to write the changed presentation to")

    ns = p.parse_args(sys.argv[1:])

    if len(ns.matches) != len(ns.replacements):
        print("There must be as many match-strings (-m) as there are replacement-strings (-r)",file=sys.stderr)
        sys.exit(1)

    print(ns)

    replacer = python_pptx_text_replacer(ns.input)
    replacements = []
    for m in range(0,len(ns.matches)):
        replacements.append( ( ns.matches[m], ns.replacements[m] ) )
    replacer.replace_text(replacements)
    replacer.write_presentation_to_file(ns.output)
    sys.exit(0)
 

